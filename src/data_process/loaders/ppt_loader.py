# -*- coding: utf-8 -*-
"""
PPT (.pptx) 文档加载器

支持文本、表格提取和幻灯片嵌入图片的 OCR。
"""

import logging
from io import BytesIO
from typing import List, Optional

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.shapes.base import BaseShape
from langchain_core.documents import Document

from src.data_process.loaders.base import BaseLoader

logger = logging.getLogger(__name__)


class PPTLoader(BaseLoader):
    """PPT 文档加载器

    支持 .pptx 格式，提取文本、表格和嵌入图片的 OCR。
    """

    SHAPE_TYPE_IMAGE = 13
    SHAPE_TYPE_GROUPED = 6

    def __init__(
        self,
        file_path: str,
        encoding: Optional[str] = None,
        ocr_enabled: bool = True,
        use_cuda: bool = True,
        **kwargs,
    ) -> None:
        """初始化 PPT 加载器

        Args:
            file_path: PPT 文件路径
            encoding: 编码（pptx 为二进制格式）
            ocr_enabled: 是否启用图片 OCR
            use_cuda: OCR 是否使用 GPU
            **kwargs: 其他参数
        """
        super().__init__(file_path, encoding, **kwargs)
        self.ocr_enabled = ocr_enabled
        self.use_cuda = use_cuda
        self._ocr_engine = None

    @property
    def ocr_engine(self):
        """延迟初始化 OCR 引擎"""
        if self._ocr_engine is None:
            try:
                from rapidocr_paddle import RapidOCR
                self._ocr_engine = RapidOCR(
                    det_use_cuda=self.use_cuda,
                    cls_use_cuda=self.use_cuda,
                    rec_use_cuda=self.use_cuda,
                )
            except ImportError:
                from rapidocr_onnxruntime import RapidOCR
                self._ocr_engine = RapidOCR()
        return self._ocr_engine

    def load(self) -> List[Document]:
        """加载 PPT 文档

        Returns:
            List[Document]: 文档列表
        """
        prs = Presentation(str(self.file_path))
        slide_contents = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = [f"[Slide {slide_num}]"]

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (x.top, x.left)
            )

            for shape in sorted_shapes:
                shape_text = self._extract_shape_content(shape)
                if shape_text:
                    slide_text.append(shape_text)

            slide_contents.append("\n".join(slide_text))

        full_text = "\n\n".join(filter(None, slide_contents))

        metadata = {
            "source": str(self.file_path),
            "file_name": self.file_path.name,
            "file_type": "pptx",
            "total_slides": len(prs.slides),
        }

        return [Document(page_content=full_text, metadata=metadata)]

    def _extract_shape_content(self, shape: BaseShape) -> str:
        """提取 Shape 内容（文本、表格、图片）"""
        parts = []

        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

        if shape.has_table:
            table_text = self._extract_table_text(shape.table)
            if table_text:
                parts.append(table_text)

        if shape.shape_type == self.SHAPE_TYPE_IMAGE:
            ocr_text = self._extract_image_ocr(shape)
            if ocr_text:
                parts.append(ocr_text)

        elif shape.shape_type == self.SHAPE_TYPE_GROUPED:
            grouped_text = self._extract_grouped_shape(shape)
            if grouped_text:
                parts.append(grouped_text)

        return "\n".join(parts)

    def _extract_table_text(self, table) -> str:
        """提取表格文本"""
        rows_text = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = " ".join(p.text.strip() for p in cell.text_frame.paragraphs)
                cells.append(cell_text)
            rows_text.append(" | ".join(cells))
        return "\n".join(rows_text)

    def _extract_image_ocr(self, shape: BaseShape) -> str:
        """对嵌入图片进行 OCR"""
        if not self.ocr_enabled:
            return ""

        try:
            img_bytes = shape.image.blob
            img = Image.open(BytesIO(img_bytes))
            img_array = np.array(img)

            result, _ = self.ocr_engine(img_array)
            if result:
                ocr_lines = [line[1] for line in result]
                return "\n".join(ocr_lines)
        except Exception as e:
            logger.debug(f"PPT 内嵌图片 OCR 失败: {e}")

        return ""

    def _extract_grouped_shape(self, shape: BaseShape) -> str:
        """递归提取组合 Shape 内容"""
        parts = []
        try:
            for child_shape in shape.shapes:
                child_text = self._extract_shape_content(child_shape)
                if child_text:
                    parts.append(child_text)
        except Exception as e:
            logger.debug(f"组合 Shape 内容提取失败: {e}")

        return "\n".join(parts)
